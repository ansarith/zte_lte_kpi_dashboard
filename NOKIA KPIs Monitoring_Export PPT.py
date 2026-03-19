import os, sys, io
import streamlit as st
import pandas as pd
import plotly.graph_objects as go
import plotly.express as px
from pptx import Presentation
from pptx.util import Inches
import matplotlib.pyplot as plt

st.set_page_config(layout="wide")

# ----------------------------------------------------
st.write("🚀 Running file:", os.path.abspath(__file__))
st.write("🟢 Python executable:", sys.executable)
# ----------------------------------------------------

# ---------------- LOAD DATA ----------------
@st.cache_data
def load_data(path):
    df = pd.read_excel(path)
    df["Begin Time"] = pd.to_datetime(df["Begin Time"], errors="coerce")

    percentage_kpis = [col for col in df.columns if "%" in col or "Rate" in col]

    for col in percentage_kpis:
        if col in df.columns and pd.api.types.is_numeric_dtype(df[col]):
            if df[col].max() <= 1.0:
                df[col] = df[col] * 100

    return df

DATA_PATH = "Performance Management-History Query-2025_LTE_KPI-Sarith-20260319094859.xlsx"
df = load_data(DATA_PATH)

st.title("📊 LTE KPI Dashboard")

# ---------------- KPI SELECTION ----------------
kpi_columns = [col for col in df.columns if col not in ["Begin Time","ENBFunction Name","Cell Name"]]
selected_kpis = st.multiselect(
    "Select KPI(s)",
    options=kpi_columns,
    default=kpi_columns[:4]
)

# 🎯 Ensure consistent order (fix color shifting)
selected_kpis = sorted(selected_kpis)

# ---------------- SITE FILTER ----------------
enodeb_selected = st.multiselect("Select ENBFunction Name", options=sorted(df["ENBFunction Name"].unique()))

# ---------------- CELL FILTER ----------------
if enodeb_selected:
    cell_options = sorted(df[df["ENBFunction Name"].isin(enodeb_selected)]["Cell Name"].unique())
else:
    cell_options = sorted(df["Cell Name"].unique())

cell_selected = st.multiselect("Select Cell Name", options=cell_options)

# ---------------- OPTIONS ----------------
daily_option = st.checkbox("📅 Daily Aggregation")
group_option = st.checkbox("🏙️ Group by Site")

# ---------------- FILTER DATAFRAME ----------------
plot_df = df.copy()

if enodeb_selected:
    plot_df = plot_df[plot_df["ENBFunction Name"].isin(enodeb_selected)]

if cell_selected:
    plot_df = plot_df[plot_df["Cell Name"].isin(cell_selected)]

# ---------------- AGGREGATION ----------------
def aggregate_data(df, kpis, daily=False, group=False):

    for kpi in kpis:
        df[kpi] = pd.to_numeric(df[kpi], errors="coerce")

    agg_dict = {}

    for kpi in kpis:
        if kpi in [
            "DL Data Total Volume (Gbyte)",
            "UL Data Total Volume (Gbyte)",
            "Total Data Total Volume (Gbyte)",
            "Ave RRC Connected Ue",
            "Max RRC Connected Ue"
        ]:
            agg_dict[kpi] = "sum"
        else:
            agg_dict[kpi] = "mean"

    if daily:
        df["Date"] = df["Begin Time"].dt.normalize()
    
        first_day = df["Date"].min()
        last_day = df["Date"].max()
    
        df = df[(df["Date"] > first_day) & (df["Date"] < last_day)]
    
        time_col = "Date"
    else:
        time_col = "Begin Time"

    if not group:
        group_cols = [time_col]

        if "Cell Name" in df.columns:
            group_cols.append("Cell Name")

        grouped = df.groupby(group_cols, as_index=False).agg(agg_dict)

    else:
        grouped = df.groupby([time_col], as_index=False).agg(agg_dict)

    return grouped

plot_df = aggregate_data(plot_df, selected_kpis, daily_option, group_option)

time_col = "Date" if daily_option else "Begin Time"

plot_df[time_col] = pd.to_datetime(plot_df[time_col], errors="coerce")
plot_df = plot_df.dropna(subset=[time_col])

plot_df["Time_str"] = plot_df[time_col].dt.strftime(
    "%Y-%m-%d" if daily_option else "%Y-%m-%d %H:%M"
)

# ---------------- DASHBOARD (PLOTLY) ----------------
# ---------------- DASHBOARD (PLOTLY) ----------------
figures_png = []

if not plot_df.empty:

    colors = px.colors.qualitative.Dark24

    # 🎯 KPI color map (ALWAYS needed)
    kpi_color_map = {
        kpi: colors[i % len(colors)]
        for i, kpi in enumerate(sorted(selected_kpis))
    }

    # 🎯 Cell color map (ONLY for non-group mode)
    if not group_option and "Cell Name" in plot_df.columns:
        unique_cells = sorted(plot_df["Cell Name"].unique())
        color_map = {
            cell: colors[i % len(colors)]
            for i, cell in enumerate(unique_cells)
        }

    cols = st.columns(2)

    for idx, selected_kpi in enumerate(selected_kpis[:4]):

        fig = go.Figure()

        # ---------- PLOTLY DASHBOARD GRAPH ----------
        if not group_option and "Cell Name" in plot_df.columns:

            for cell in sorted(plot_df["Cell Name"].unique()):

                cell_df = plot_df[plot_df["Cell Name"] == cell]

                fig.add_trace(
                    go.Scatter(
                        x=cell_df["Time_str"],
                        y=cell_df[selected_kpi],
                        mode="lines+markers",
                        name=cell,
                        line=dict(color=color_map[cell])
                    )
                )

        else:

            fig.add_trace(
                go.Scatter(
                    x=plot_df["Time_str"],
                    y=plot_df[selected_kpi],
                    mode="lines+markers",
                    name=selected_kpi,
                    line=dict(color=kpi_color_map[selected_kpi])
                )
            )

        fig.update_layout(
            height=420,
            width=900,
            title=dict(text=selected_kpi, x=0.5),
            hovermode="x unified",
            margin=dict(l=40, r=120, t=60, b=40)
        )

        # Show in Streamlit
        cols[idx % 2].plotly_chart(fig, use_container_width=False)

        # ---------- EXPORT SAME FIGURE TO PNG ----------
        img_bytes = fig.to_image(
            format="png",
            width=900,
            height=420,
            scale=2
        )

        buf = io.BytesIO(img_bytes)
        figures_png.append(buf)

# ---------------- CREATE PPT ----------------
def create_ppt(figures_png):

    prs = Presentation()

    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    positions = [
        (Inches(0.5), Inches(0.5)),
        (Inches(6.9), Inches(0.5)),
        (Inches(0.5), Inches(4.0)),
        (Inches(6.9), Inches(4.0))
    ]

    chart_width = Inches(6.08)
    chart_height = Inches(3.04)

    for idx, buf in enumerate(figures_png):

        if idx % 4 == 0:
            slide = prs.slides.add_slide(prs.slide_layouts[5])

        pos_idx = idx % 4

        slide.shapes.add_picture(
            buf,
            positions[pos_idx][0],
            positions[pos_idx][1],
            width=chart_width,
            height=chart_height
        )

    ppt_buffer = io.BytesIO()

    prs.save(ppt_buffer)

    ppt_buffer.seek(0)

    return ppt_buffer

# ---------------- DOWNLOAD PPT ----------------
if figures_png:

    ppt_file = create_ppt(figures_png)

    st.download_button(
        "📊 Download PowerPoint Report",
        data=ppt_file,
        file_name="LTE_KPI_Report.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )

else:

    st.warning("⚠️ No data available for the selected filters.")




